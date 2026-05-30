"""
生成答辩 PPT 脚本
运行：python create_ppt.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# =========== 颜色常量 ===========
COLOR_PRIMARY = RGBColor(0x1E, 0x1B, 0x4B)     # 深蓝紫色
COLOR_ACCENT = RGBColor(0x4F, 0x46, 0xE5)       # 靛蓝
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_BLACK = RGBColor(0x33, 0x33, 0x33)
COLOR_GRAY = RGBColor(0x66, 0x66, 0x66)
COLOR_LIGHT_BG = RGBColor(0xF5, 0xF5, 0xF5)
COLOR_GREEN = RGBColor(0x10, 0xB9, 0x81)
COLOR_RED = RGBColor(0xEF, 0x44, 0x44)
COLOR_ORANGE = RGBColor(0xE6, 0xA2, 0x3C)


def add_bg(slide, color=COLOR_PRIMARY):
    """给幻灯片设置纯色背景"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_bg(slide, color=COLOR_ACCENT, left=0, top=0, width=None, height=None):
    """添加矩形色块"""
    if width is None:
        width = prs.slide_width
    if height is None:
        height = prs.slide_height
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text_box(slide, text, left, top, width, height, font_size=18, bold=False,
                 color=COLOR_BLACK, alignment=PP_ALIGN.LEFT, font_name="微软雅黑"):
    """添加文本框"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_slide(slide, title, items, title_color=COLOR_WHITE, text_color=COLOR_WHITE):
    """添加带要点的幻灯片"""
    add_bg(slide, COLOR_PRIMARY)
    add_text_box(slide, title, Inches(0.8), Inches(0.5), Inches(11.5), Inches(1),
                 font_size=36, bold=True, color=title_color)
    # 装饰线
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.4), Inches(2), Pt(4))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_ACCENT
    shape.line.fill.background()
    # 要点列表
    y = Inches(1.8)
    for item in items:
        add_text_box(slide, f"  ▸  {item}", Inches(1.2), y, Inches(10.5), Inches(0.6),
                     font_size=20, color=text_color)
        y += Inches(0.55)


# =========== Slide 1: 封面 ===========
slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
add_bg(slide, COLOR_PRIMARY)

# 主标题
add_text_box(slide, "个人收支财务管理系统", Inches(1.5), Inches(1.8), Inches(10), Inches(1.5),
             font_size=48, bold=True, color=COLOR_WHITE, alignment=PP_ALIGN.CENTER)

# 装饰线
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.5), Inches(3.5), Inches(4), Pt(4))
shape.fill.solid()
shape.fill.fore_color.rgb = COLOR_ACCENT
shape.line.fill.background()

# 副标题
add_text_box(slide, "软件工程课程设计  |  答辩演示", Inches(1.5), Inches(3.8), Inches(10), Inches(0.8),
             font_size=24, color=RGBColor(0xA5, 0xB4, 0xFC), alignment=PP_ALIGN.CENTER)

# 团队信息
info = "小组成员：组长（总统筹岗）· 需求设计岗 · 核心编码岗 · 界面数据岗 · 测试答辩岗"
add_text_box(slide, info, Inches(1.5), Inches(5.2), Inches(10), Inches(0.6),
             font_size=18, color=RGBColor(0xC7, 0xD2, 0xFE), alignment=PP_ALIGN.CENTER)

add_text_box(slide, "2026年5月", Inches(1.5), Inches(5.8), Inches(10), Inches(0.5),
             font_size=16, color=RGBColor(0x94, 0xA3, 0xB8), alignment=PP_ALIGN.CENTER)


# =========== Slide 2: 目录 ===========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, COLOR_PRIMARY)
add_text_box(slide, "目  录", Inches(0.8), Inches(0.5), Inches(5), Inches(1),
             font_size=40, bold=True, color=COLOR_WHITE)

items = [
    "项目背景与目标",
    "技术架构",
    "功能模块展示",
    "数据库设计",
    "测试总结",
    "团队分工与总结"
]
y = Inches(1.6)
for i, item in enumerate(items, 1):
    add_text_box(slide, f"  {i:02d}    {item}", Inches(1.5), y, Inches(10), Inches(0.7),
                 font_size=24, color=COLOR_WHITE)
    y += Inches(0.75)


# =========== Slide 3: 项目背景 ===========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bullet_slide(slide, "一、项目背景与目标", [
    "痛点：传统手工记账易遗漏、分类不清晰、分析不直观、多账户管理困难",
    "目标用户：学生、职场新人、家庭用户",
    "核心需求：多账户管理、收支记录、分类管理、预算管控、统计分析",
    "开发模式：前后端分离，团队 5 人协作完成",
])


# =========== Slide 4: 技术架构 ===========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, COLOR_PRIMARY)
add_text_box(slide, "二、技术架构", Inches(0.8), Inches(0.5), Inches(11.5), Inches(1),
             font_size=36, bold=True, color=COLOR_WHITE)
# 装饰线
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.4), Inches(2), Pt(4))
shape.fill.solid()
shape.fill.fore_color.rgb = COLOR_ACCENT
shape.line.fill.background()

# 架构图 - 前端
box_w = Inches(3.8)
box_h = Inches(1.2)
front_rect = add_shape_bg(slide, RGBColor(0x4F, 0x46, 0xE5), Inches(0.8), Inches(1.8), box_w, box_h)
add_text_box(slide, "前端（Vue 3 SPA）", Inches(1.0), Inches(1.95), box_w, Inches(0.5),
             font_size=22, bold=True, color=COLOR_WHITE, alignment=PP_ALIGN.CENTER)
add_text_box(slide, "Element Plus · Vue Router · Axios", Inches(1.0), Inches(2.4), box_w, Inches(0.4),
             font_size=14, color=RGBColor(0xC7, 0xD2, 0xFE), alignment=PP_ALIGN.CENTER)

# 中间箭头
add_text_box(slide, "⟷  HTTP JSON (RESTful API)  ⟷", Inches(4.8), Inches(2.1), Inches(3.5), Inches(0.6),
             font_size=18, color=COLOR_WHITE, alignment=PP_ALIGN.CENTER)

# 架构图 - 后端
back_rect = add_shape_bg(slide, RGBColor(0x05, 0x91, 0x69), Inches(8.5), Inches(1.8), box_w, box_h)
add_text_box(slide, "后端（Flask 3.0）", Inches(8.7), Inches(1.95), box_w, Inches(0.5),
             font_size=22, bold=True, color=COLOR_WHITE, alignment=PP_ALIGN.CENTER)
add_text_box(slide, "Route → Service → DAO", Inches(8.7), Inches(2.4), box_w, Inches(0.4),
             font_size=14, color=RGBColor(0xA7, 0xF3, 0xD0), alignment=PP_ALIGN.CENTER)

# 数据库层
db_rect = add_shape_bg(slide, RGBColor(0xD9, 0x77, 0x06), Inches(8.5), Inches(3.5), box_w, box_h)
add_text_box(slide, "MySQL 数据库", Inches(8.7), Inches(3.65), box_w, Inches(0.5),
             font_size=22, bold=True, color=COLOR_WHITE, alignment=PP_ALIGN.CENTER)
add_text_box(slide, "5 张表 · InnoDB · utf8mb4", Inches(8.7), Inches(4.1), box_w, Inches(0.4),
             font_size=14, color=RGBColor(0xFE, 0xEB, 0xC8), alignment=PP_ALIGN.CENTER)

# 向下箭头
add_text_box(slide, "↓ SQL", Inches(9.5), Inches(3.1), Inches(1.5), Inches(0.4),
             font_size=16, color=COLOR_WHITE, alignment=PP_ALIGN.CENTER)

# 部署说明
items = [
    "后端启动：python app.py（端口 5000）",
    "前端启动：npm run dev（端口 5173，代理 /api → :5000）",
    "环境要求：Python 3.10+ · Node 18+ · MySQL 5.7+"
]
y = Inches(5.0)
for item in items:
    add_text_box(slide, f"  ●  {item}", Inches(0.8), y, Inches(11.5), Inches(0.45),
                 font_size=16, color=RGBColor(0xE2, 0xE8, 0xF0))
    y += Inches(0.45)


# =========== Slide 5: 功能模块 ===========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, COLOR_PRIMARY)
add_text_box(slide, "三、功能模块展示", Inches(0.8), Inches(0.5), Inches(11.5), Inches(1),
             font_size=36, bold=True, color=COLOR_WHITE)
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.4), Inches(2), Pt(4))
shape.fill.solid()
shape.fill.fore_color.rgb = COLOR_ACCENT
shape.line.fill.background()

modules = [
    ("📊 财务总览", "总收入/支出/结余概览，最近交易列表", COLOR_GREEN),
    ("🏦 账户管理", "多账户增删改查，支持现金/银行卡/微信/支付宝", COLOR_ACCENT),
    ("💰 收支记录", "按账户/日期筛选，添加收支，账户间转账，CSV导出", RGBColor(0xE6, 0xA2, 0x3C)),
    ("🏷️ 分类管理", "自定义收入/支出两级分类，预置 33 条默认分类", RGBColor(0x05, 0x91, 0x69)),
    ("🎯 预算管理", "按月设置预算，实时计算剩余，超支 80% 预警", RGBColor(0xD9, 0x77, 0x06)),
    ("📈 统计分析", "多维度统计，按日/周/月/年趋势分析", RGBColor(0x8B, 0x5C, 0xF6)),
]

y = Inches(1.7)
for title, desc, color in modules:
    rect = add_shape_bg(slide, RGBColor(0x31, 0x2E, 0x81), Inches(0.8), y, Inches(11.2), Inches(0.75))
    add_text_box(slide, title, Inches(1.0), y + Inches(0.08), Inches(2.5), Inches(0.5),
                 font_size=20, bold=True, color=color)
    add_text_box(slide, desc, Inches(3.5), y + Inches(0.08), Inches(8.5), Inches(0.5),
                 font_size=18, color=RGBColor(0xE2, 0xE8, 0xF0))
    y += Inches(0.85)


# =========== Slide 6: 数据库设计 ===========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, COLOR_PRIMARY)
add_text_box(slide, "四、数据库设计", Inches(0.8), Inches(0.5), Inches(11.5), Inches(1),
             font_size=36, bold=True, color=COLOR_WHITE)
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.4), Inches(2), Pt(4))
shape.fill.solid()
shape.fill.fore_color.rgb = COLOR_ACCENT
shape.line.fill.background()

# 表列表
tables = [
    ("accounts（账户表）",       "id, name, type, balance, created_at"),
    ("transactions（交易表）",   "id, account_id(FK), type, category, amount, date"),
    ("categories（分类表）",     "id, type, main, sub（预置33条）"),
    ("budgets（预算表）",        "id, year, month, category, amount"),
    ("logs（日志表）",           "id, time, user, action, detail"),
]

y = Inches(1.7)
for name, fields in tables:
    add_text_box(slide, f"  ●  {name}", Inches(1.2), y, Inches(5), Inches(0.45),
                 font_size=20, bold=True, color=COLOR_ACCENT)
    add_text_box(slide, fields, Inches(1.5), y + Inches(0.4), Inches(9), Inches(0.4),
                 font_size=16, color=RGBColor(0xE2, 0xE8, 0xF0))
    y += Inches(0.85)

add_text_box(slide, "关键技术：InnoDB 引擎 · 外键级联删除 · DBUtils 连接池 · 事务上下文管理",
             Inches(0.8), Inches(6.2), Inches(11.5), Inches(0.5),
             font_size=16, color=RGBColor(0x94, 0xA3, 0xB8))


# =========== Slide 7: 测试总结 ===========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, COLOR_PRIMARY)
add_text_box(slide, "五、测试总结", Inches(0.8), Inches(0.5), Inches(11.5), Inches(1),
             font_size=36, bold=True, color=COLOR_WHITE)
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.4), Inches(2), Pt(4))
shape.fill.solid()
shape.fill.fore_color.rgb = COLOR_ACCENT
shape.line.fill.background()

# 测试统计
stats = [
    ("测试类型",   "用例数", "通过率"),
    ("单元测试",   "39", "100%"),
    ("系统测试",   "9",  "100%"),
    ("合计",       "48", "100%"),
]
y = Inches(1.8)
col_x = [Inches(1.0), Inches(5.0), Inches(8.5)]
for row in stats:
    for i, val in enumerate(row):
        is_header = row == stats[0]
        c = RGBColor(0xA5, 0xB4, 0xFC) if is_header else COLOR_WHITE
        fs = 18 if is_header else 20
        b = is_header
        add_text_box(slide, val, col_x[i], y, Inches(3), Inches(0.5),
                     font_size=fs, bold=b, color=c, alignment=PP_ALIGN.CENTER)
    y += Inches(0.5)

# 底部总结
add_text_box(slide, "✅ 全部 48 个测试用例通过，通过率 100%", Inches(0.8), Inches(4.2), Inches(11.5), Inches(0.5),
             font_size=22, bold=True, color=COLOR_GREEN, alignment=PP_ALIGN.CENTER)
add_text_box(slide, "系统运行稳定，功能完整，异常处理完善，建议验收通过",
             Inches(0.8), Inches(4.8), Inches(11.5), Inches(0.5),
             font_size=18, color=RGBColor(0xE2, 0xE8, 0xF0), alignment=PP_ALIGN.CENTER)

# 测试模块
test_items = "测试范围：账户管理 · 收支记录 · 转账 · 分类管理 · 预算管理 · 统计分析 · 操作日志"
add_text_box(slide, test_items, Inches(0.8), Inches(5.6), Inches(11.5), Inches(0.5),
             font_size=16, color=RGBColor(0x94, 0xA3, 0xB8), alignment=PP_ALIGN.CENTER)


# =========== Slide 8: 团队分工 ===========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, COLOR_PRIMARY)
add_text_box(slide, "六、团队分工与总结", Inches(0.8), Inches(0.5), Inches(11.5), Inches(1),
             font_size=36, bold=True, color=COLOR_WHITE)
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.4), Inches(2), Pt(4))
shape.fill.solid()
shape.fill.fore_color.rgb = COLOR_ACCENT
shape.line.fill.background()

roles = [
    ("👑 组长（总统筹岗）", "可行性分析、进度把控、问题协调、文档框架"),
    ("📋 需求设计岗",       "需求分析、总体设计、详细设计、界面草图"),
    ("💻 核心编码岗",       "核心业务编码、代码调试、编码规范制定"),
    ("🎨 界面数据岗",       "界面开发、数据库实现、截图采集"),
    ("🧪 测试答辩岗",       "软件测试、文档汇总、答辩PPT准备"),
]

y = Inches(1.6)
for role, duty in roles:
    rect = add_shape_bg(slide, RGBColor(0x31, 0x2E, 0x81), Inches(0.8), y, Inches(11.2), Inches(0.85))
    add_text_box(slide, role, Inches(1.0), y + Inches(0.08), Inches(4), Inches(0.5),
                 font_size=20, bold=True, color=COLOR_ACCENT)
    add_text_box(slide, duty, Inches(5.0), y + Inches(0.08), Inches(7), Inches(0.5),
                 font_size=18, color=RGBColor(0xE2, 0xE8, 0xF0))
    y += Inches(0.95)


# =========== Slide 9: 致谢 ===========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, COLOR_PRIMARY)

add_text_box(slide, "感谢聆听", Inches(1.5), Inches(2.5), Inches(10), Inches(1.2),
             font_size=52, bold=True, color=COLOR_WHITE, alignment=PP_ALIGN.CENTER)

shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5), Inches(3.8), Inches(3), Pt(4))
shape.fill.solid()
shape.fill.fore_color.rgb = COLOR_ACCENT
shape.line.fill.background()

add_text_box(slide, "个人收支财务管理系统", Inches(1.5), Inches(4.2), Inches(10), Inches(0.6),
             font_size=22, color=RGBColor(0xA5, 0xB4, 0xFC), alignment=PP_ALIGN.CENTER)

add_text_box(slide, "有任何问题，欢迎提出！", Inches(1.5), Inches(5.2), Inches(10), Inches(0.5),
             font_size=18, color=RGBColor(0x94, 0xA3, 0xB8), alignment=PP_ALIGN.CENTER)


# 保存
output_path = "E:\\University\\software_project\\documents\\答辩PPT.pptx"
prs.save(output_path)
print(f"PPT 已保存到: {output_path}")
